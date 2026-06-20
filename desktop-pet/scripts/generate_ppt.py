import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    COLOR_PRIMARY = RGBColor(31, 122, 90)     # Forest Green (#1f7a5a)
    COLOR_BG = RGBColor(246, 247, 242)       # Light Off-White/Gray (#f6f7f2)
    COLOR_TEXT_DARK = RGBColor(23, 33, 27)    # Dark Slate (#17211b)
    COLOR_TEXT_MUTED = RGBColor(102, 116, 109) # Soft Gray (#66746d)
    COLOR_WHITE = RGBColor(255, 255, 255)     # Card background
    COLOR_BORDER = RGBColor(210, 218, 214)    # Thin border
    COLOR_ACCENT = RGBColor(218, 145, 0)      # Warm gold/orange for highlight (#da9100)

    # Fonts
    FONT_TITLE = "Microsoft YaHei"
    FONT_BODY = "Microsoft YaHei"
    
    # Helper to set background
    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw a rectangle (card)
    def draw_card(slide, left, top, width, height, fill_color, border_color=None, line_width=1):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()
        return shape

    # Helper to draw a line
    def draw_line(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # Helper to add a textbox
    def add_text(slide, text, left, top, width, height, font_size, font_color, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_BODY, word_wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        return txBox

    # Helper to add multi-paragraph text block
    def add_rich_text(slide, left, top, width, height, paragraphs_data, font_name=FONT_BODY):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        for i, p_data in enumerate(paragraphs_data):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = p_data.get("text", "")
            p.alignment = p_data.get("align", PP_ALIGN.LEFT)
            p.font.name = font_name
            p.font.size = p_data.get("size", Pt(12))
            p.font.color.rgb = p_data.get("color", COLOR_TEXT_DARK)
            p.font.bold = p_data.get("bold", False)
            
            if "space_after" in p_data:
                p.space_after = p_data["space_after"]
            if "space_before" in p_data:
                p.space_before = p_data["space_before"]
            if "level" in p_data:
                p.level = p_data["level"]
        return txBox

    # Helper to add standard header and footer for content slides
    def add_header_footer(slide, title_text, slide_index):
        # Header background banner
        draw_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), COLOR_WHITE)
        draw_line(slide, Inches(0), Inches(1.18), Inches(13.333), Inches(0.02), COLOR_PRIMARY)
        
        # Header Title
        add_text(slide, title_text, Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.6), Pt(24), COLOR_PRIMARY, bold=True, font_name=FONT_TITLE)
        
        # Footer
        draw_line(slide, Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.01), COLOR_BORDER)
        add_text(slide, "《灵契·伴读灵》-- 词汇羁绊养成系统 | 课程大作业答辩汇报", Inches(0.8), Inches(7.0), Inches(8.0), Inches(0.3), Pt(10), COLOR_TEXT_MUTED)
        add_text(slide, f"{slide_index} / 8", Inches(11.533), Inches(7.0), Inches(1.0), Inches(0.3), Pt(10), COLOR_TEXT_MUTED, align=PP_ALIGN.RIGHT)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Cover)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide1, COLOR_BG)
    
    # Left decorative green panel
    draw_card(slide1, Inches(0), Inches(0), Inches(4.8), Inches(7.5), COLOR_PRIMARY)
    
    # Large watermark / English text in left panel
    add_text(slide1, "WORD TRAINER\n& DETAILED\nCOMPANION", Inches(0.5), Inches(0.8), Inches(3.8), Inches(2.0), Pt(32), RGBColor(45, 145, 110), bold=True, word_wrap=True)
    add_text(slide1, "C++ & WebView2 Hybrid Desktop Application", Inches(0.5), Inches(2.8), Inches(3.8), Inches(1.0), Pt(11), RGBColor(165, 220, 195), bold=False, word_wrap=True)
    
    # Logo text / Brand in left panel bottom
    add_text(slide1, "THE BOND SYSTEM", Inches(0.5), Inches(6.4), Inches(3.8), Inches(0.5), Pt(12), RGBColor(165, 220, 195), bold=True)
    
    # Right Main content area (Start at x=5.4)
    # Project Title
    add_text(slide1, "《灵契·伴读灵》", Inches(5.4), Inches(1.6), Inches(7.0), Inches(0.9), Pt(42), COLOR_PRIMARY, bold=True, font_name=FONT_TITLE)
    add_text(slide1, "—— 词汇羁绊养成系统答辩汇报", Inches(5.4), Inches(2.5), Inches(7.0), Inches(0.6), Pt(22), COLOR_TEXT_DARK, bold=True, font_name=FONT_TITLE)
    
    # Title decorative accent line
    draw_line(slide1, Inches(5.4), Inches(3.3), Inches(4.5), Inches(0.04), COLOR_PRIMARY)
    
    # Subtitle or description
    add_text(slide1, "基于 C++ 与 WebView2 混合架构的系统级单词背诵与桌宠伴学应用", Inches(5.4), Inches(3.6), Inches(7.0), Inches(0.5), Pt(13), COLOR_TEXT_MUTED)
    
    # Details Box (Card)
    draw_card(slide1, Inches(5.4), Inches(4.4), Inches(4.8), Inches(1.8), COLOR_WHITE, COLOR_BORDER)
    details = [
        {"text": "报告人 :  TYH", "size": Pt(12), "color": COLOR_TEXT_DARK, "space_after": Pt(6), "bold": True},
        {"text": "开发语言 :  C++ & HTML5 / JavaScript", "size": Pt(11), "color": COLOR_TEXT_DARK, "space_after": Pt(6)},
        {"text": "核心架构 :  Win32 + Direct2D + SAPI + WebView2", "size": Pt(11), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide1, Inches(5.7), Inches(4.7), Inches(4.2), Inches(1.2), details)
    
    # Insert Pet Image on the boundary / in the right-bottom area
    pet_img_path = r"c:\Users\TYH\Documents\Codex\2026-06-11\c\outputs\word_trainer_web\desktop-pet\res\aimee_pet_hd.png"
    if os.path.exists(pet_img_path):
        # We specify width only to preserve aspect ratio
        slide1.shapes.add_picture(pet_img_path, Inches(10.0), Inches(3.5), width=Inches(2.8))

    # -------------------------------------------------------------
    # SLIDE 2: Problems & Goals (问题与目标)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2, COLOR_BG)
    add_header_footer(slide2, "01 / 背景痛点与项目设计目标", 2)
    
    # Left Card: Traditional Problems
    draw_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.9), COLOR_WHITE, COLOR_BORDER)
    # Title badge for Problems
    draw_card(slide2, Inches(0.8), Inches(1.6), Inches(2.2), Inches(0.5), RGBColor(230, 90, 90)) # Red accent for problems
    add_text(slide2, "传统背单词软件痛点", Inches(1.0), Inches(1.7), Inches(2.0), Inches(0.3), Pt(12), COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    problems_data = [
        {"text": "• 学习交互枯燥单调", "size": Pt(15), "color": COLOR_TEXT_DARK, "bold": True, "space_before": Pt(14), "space_after": Pt(4)},
        {"text": "   单纯拼写与选择题，缺乏趣味与伴学机制，用户难以长期坚持。", "size": Pt(11), "color": COLOR_TEXT_MUTED, "space_after": Pt(10)},
        
        {"text": "• 计划与复习提醒微弱", "size": Pt(15), "color": COLOR_TEXT_DARK, "bold": True, "space_after": Pt(4)},
        {"text": "   提醒仅停留在应用内部，软件关闭后失去触达，导致复习极易中断。", "size": Pt(11), "color": COLOR_TEXT_MUTED, "space_after": Pt(10)},
        
        {"text": "• 语音朗读引起界面卡顿", "size": Pt(15), "color": COLOR_TEXT_DARK, "bold": True, "space_after": Pt(4)},
        {"text": "   常规网页调用TTS发音时，易因网络延迟或单线程阻塞导致界面短时间卡死。", "size": Pt(11), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide2, Inches(1.1), Inches(2.3), Inches(5.0), Inches(4.0), problems_data)
    
    # Right Card: Design Goals
    draw_card(slide2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.9), COLOR_WHITE, COLOR_BORDER)
    # Title badge for Goals
    draw_card(slide2, Inches(6.9), Inches(1.6), Inches(2.2), Inches(0.5), COLOR_PRIMARY) # Green accent for goals
    add_text(slide2, "《灵契》项目设计目标", Inches(7.1), Inches(1.7), Inches(2.0), Inches(0.3), Pt(12), COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    goals_data = [
        {"text": "• 引入桌面“伴读灵”角色", "size": Pt(15), "color": COLOR_PRIMARY, "bold": True, "space_before": Pt(14), "space_after": Pt(4)},
        {"text": "   将单词记忆转化为人机羁绊互动，融入日常桌面使用，提升成就感。", "size": Pt(11), "color": COLOR_TEXT_MUTED, "space_after": Pt(10)},
        
        {"text": "• 系统级气泡常驻提醒", "size": Pt(15), "color": COLOR_PRIMARY, "bold": True, "space_after": Pt(4)},
        {"text": "   利用 Win32 透明悬浮窗实现常驻提醒，支持轮询错题与待复习强通知。", "size": Pt(11), "color": COLOR_TEXT_MUTED, "space_after": Pt(10)},
        
        {"text": "• 多线程非阻塞 TTS 发音", "size": Pt(15), "color": COLOR_PRIMARY, "bold": True, "space_after": Pt(4)},
        {"text": "   C++ 后台线程独立加载 Windows SAPI 引擎朗读单词，前端交互极致流畅。", "size": Pt(11), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide2, Inches(7.2), Inches(2.3), Inches(5.0), Inches(4.0), goals_data)

    # -------------------------------------------------------------
    # SLIDE 3: Overall Architecture (总体架构)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3, COLOR_BG)
    add_header_footer(slide3, "02 / 混合架构与模块划分", 3)
    
    # Main subtitle
    add_text(slide3, "项目以 DesktopPet.exe 为核心载体，构建了 WebView2 界面、C++ 核心与悬浮桌宠三位一体的混合客户端架构：", Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.4), Pt(12), COLOR_TEXT_DARK)
    
    col_width = Inches(3.64)
    gap = Inches(0.4)
    left_start = Inches(0.8)
    card_height = Inches(4.3)
    card_top = Inches(1.8)
    
    # Column 1: WebView2 Front-End
    draw_card(slide3, left_start, card_top, col_width, card_height, COLOR_WHITE, COLOR_PRIMARY)
    draw_card(slide3, left_start, card_top, col_width, Inches(0.6), COLOR_PRIMARY)
    add_text(slide3, "WebView2 前端交互层", left_start, card_top + Inches(0.15), col_width, Inches(0.4), Pt(13), COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    fe_items = [
        {"text": "• 核心学习界面", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK, "space_before": Pt(10)},
        {"text": "   基于 HTML5/JS 构建现代化背词交互", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(6)},
        {"text": "• 助记卡片与生词本", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK},
        {"text": "   直观呈现记忆联想与自定义生词整理", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(6)},
        {"text": "• 控制台与统计图表", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK},
        {"text": "   桌宠形象切换、日常练习进度大数据分析", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide3, left_start + Inches(0.2), card_top + Inches(0.8), col_width - Inches(0.4), card_height - Inches(1.0), fe_items)
    
    # Column 2: C++ Backend Service
    draw_card(slide3, left_start + col_width + gap, card_top, col_width, card_height, COLOR_WHITE, COLOR_PRIMARY)
    draw_card(slide3, left_start + col_width + gap, card_top, col_width, Inches(0.6), COLOR_PRIMARY)
    add_text(slide3, "C++ 核心服务层", left_start + col_width + gap, card_top + Inches(0.15), col_width, Inches(0.4), Pt(13), COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    be_items = [
        {"text": "• 本地 API 响应机制", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK, "space_before": Pt(10)},
        {"text": "   在宿主进程实现状态、生词等数据交互", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(6)},
        {"text": "• 多线程 TTS 朗读引擎", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK},
        {"text": "   多线程后台调度 Windows SAPI 朗读", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(6)},
        {"text": "• 存储与持久化管理", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK},
        {"text": "   加密读写本地 JSON 文件保存学习记录", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide3, left_start + col_width + gap + Inches(0.2), card_top + Inches(0.8), col_width - Inches(0.4), card_height - Inches(1.0), be_items)
    
    # Column 3: C++ DeskPet Render
    draw_card(slide3, left_start + (col_width + gap)*2, card_top, col_width, card_height, COLOR_WHITE, COLOR_PRIMARY)
    draw_card(slide3, left_start + (col_width + gap)*2, card_top, col_width, Inches(0.6), COLOR_PRIMARY)
    add_text(slide3, "C++ 桌宠渲染层", left_start + (col_width + gap)*2, card_top + Inches(0.15), col_width, Inches(0.4), Pt(13), COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    pet_items = [
        {"text": "• Win32 透明悬浮窗口", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK, "space_before": Pt(10)},
        {"text": "   常驻置顶，支持鼠标穿透与拖拽定位", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(6)},
        {"text": "• Direct2D / WIC 渲染", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK},
        {"text": "   硬件加速绘制高清角色图像与半透明气泡", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(6)},
        {"text": "• 气泡同步提醒与托盘", "size": Pt(13), "bold": True, "color": COLOR_TEXT_DARK},
        {"text": "   桌宠气泡与前端背词同步对话，托盘防打扰", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide3, left_start + (col_width + gap)*2 + Inches(0.2), card_top + Inches(0.8), col_width - Inches(0.4), card_height - Inches(1.0), pet_items)
    
    # Bottom Communication Bridge Info Band
    draw_card(slide3, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.5), COLOR_BG, COLOR_BORDER)
    add_text(slide3, "通信核心桥梁：前端使用 chrome.webview.postMessage 触发，C++ 拦截 WebMessageReceived 分发，效率远超传统局域网 Socket 交互。", Inches(1.0), Inches(6.32), Inches(11.333), Inches(0.3), Pt(10.5), COLOR_TEXT_DARK, bold=True)

    # -------------------------------------------------------------
    # SLIDE 4: C++ Implementation Details (C++ 核心技术实现)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide4, COLOR_BG)
    add_header_footer(slide4, "03 / C++ 核心技术实现与代码架构", 4)
    
    card_w = Inches(5.6)
    card_h = Inches(2.2)
    top_row1 = Inches(1.8)
    top_row2 = Inches(4.3)
    col1 = Inches(0.8)
    col2 = Inches(6.9)
    
    # Grid 1: Win32 消息与窗口
    draw_card(slide4, col1, top_row1, card_w, card_h, COLOR_WHITE, COLOR_BORDER)
    draw_card(slide4, col1, top_row1, Inches(0.08), card_h, COLOR_PRIMARY) # Left colored vertical accent line
    g1_text = [
        {"text": "Win32 窗口管理与系统交互", "size": Pt(14), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(6)},
        {"text": "• 悬浮窗设计: 使用 WS_EX_LAYERED 创建无边框透明窗口，通过 UpdateLayeredWindow 实现像素级 Alpha 透明通道。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(4)},
        {"text": "• 系统托盘 & 热键: 实现 Shell_NotifyIcon 托盘菜单，注册 RegisterHotKey 全局快速隐藏/唤醒桌宠与主面板。", "size": Pt(10.5), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide4, col1 + Inches(0.25), top_row1 + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3), g1_text)
    
    # Grid 2: D2D & WIC 渲染
    draw_card(slide4, col2, top_row1, card_w, card_h, COLOR_WHITE, COLOR_BORDER)
    draw_card(slide4, col2, top_row1, Inches(0.08), card_h, COLOR_PRIMARY)
    g2_text = [
        {"text": "Direct2D + WIC 高保真绘制", "size": Pt(14), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(6)},
        {"text": "• 高清角色: 弃用老旧 GDI，使用 WIC (Windows Imaging Component) 解码 PNG 透明通道，保留宠物毛发抗锯齿细节。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(4)},
        {"text": "• 对话气泡: 基于 D2D Geometry 动态绘制圆角对话框与指示小三角，支持多行文本自适应换行，彻底杜绝文本截断。", "size": Pt(10.5), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide4, col2 + Inches(0.25), top_row1 + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3), g2_text)
    
    # Grid 3: 多线程 TTS
    draw_card(slide4, col1, top_row2, card_w, card_h, COLOR_WHITE, COLOR_BORDER)
    draw_card(slide4, col1, top_row2, Inches(0.08), card_h, COLOR_PRIMARY)
    g3_text = [
        {"text": "多线程 Windows SAPI 语音朗读", "size": Pt(14), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(6)},
        {"text": "• 非阻塞发音: 初始化 COM 库与 ISpVoice 接口，采用 std::thread 创建后台工作线程异步调用 Speak() 朗读单词。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(4)},
        {"text": "• 线程安全控制: 使用互斥锁保证多线程状态同步，当用户连续点击发音时自动打断上一个发音，保障用户流畅答题。", "size": Pt(10.5), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide4, col1 + Inches(0.25), top_row2 + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3), g3_text)
    
    # Grid 4: 序列化与文件持久化
    draw_card(slide4, col2, top_row2, card_w, card_h, COLOR_WHITE, COLOR_BORDER)
    draw_card(slide4, col2, top_row2, Inches(0.08), card_h, COLOR_PRIMARY)
    g4_text = [
        {"text": "Boost.JSON 通信解析与持久化存储", "size": Pt(14), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(6)},
        {"text": "• 结构化通信: 选用高性能 Boost.JSON 作为网络/桥接传输载体，零开销解析复杂的前后端答题状态 JSON 结构。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(4)},
        {"text": "• 离线数据同步: 采用 std::fstream 加密持久化存储用户学习日志、生词数据、错题集、桌宠屏幕坐标，支持离线继续学习。", "size": Pt(10.5), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide4, col2 + Inches(0.25), top_row2 + Inches(0.15), card_w - Inches(0.4), card_h - Inches(0.3), g4_text)

    # -------------------------------------------------------------
    # SLIDE 5: WebView2 & Communication (WebView2 内嵌与通信桥接)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide5, COLOR_BG)
    add_header_footer(slide5, "04 / WebView2 桥接通信机制与 API 设计", 5)
    
    # Left: Bridge Diagram Panel
    draw_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide5, "WebView2 通信桥接机制", Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    bridge_steps = [
        {"text": "1. 统一程序入口与内嵌加载", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(10), "space_after": Pt(2)},
        {"text": "   C++ 启动时拉起 Microsoft WebView2 内核，直接加载本地包内文件。前端无需运行于公开浏览器，形成纯桌面应用形态。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "2. 单向低时延原生消息传递", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   前端直接调用 chrome.webview.postMessage(json_str)，消息在宿主内部进程级分发，彻底抛弃传统的局域网 HTTP 握手过程。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "3. 宿主端拦截分发机制", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   C++ 绑定 WebMessageReceived 监听器，接收请求后利用底层模块完成 TTS、数据落地，并通过 JS ExecuteScript 异步响应前端。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide5, Inches(1.1), Inches(2.4), Inches(5.0), Inches(3.8), bridge_steps)
    
    # Right: API Design Panel
    draw_card(slide5, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide5, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide5, "前端与 C++ 交互的核心 API 路由设计", Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    api_items = [
        {"text": "• [API]  /status  --  心跳与服务监控", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(10), "space_after": Pt(2)},
        {"text": "   校验 C++ 核心后台与 HTTP 服务是否运行正常，保障前后端通讯同步。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "• [API]  /speak  --  后台非阻塞发音", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   网页传递英文单词 JSON，由 C++ SAPI 语音线程播放，绝不干扰前端主线程渲染。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "• [API]  /record  &  /study-state  --  数据与状态同步", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   上传用户答题明细，下发最新词汇复习表、错题项。C++ 解析后立即写入物理磁盘。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "• [API]  /bubble  --  桌宠消息同步播报", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   让前端能随时唤醒桌宠气泡对话框，进行闲聊提示或错题高亮警示。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide5, Inches(7.2), Inches(2.4), Inches(5.0), Inches(3.8), api_items)

    # -------------------------------------------------------------
    # SLIDE 6: Study Algorithm & Data Flow (学习计划与数据流)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide6, COLOR_BG)
    add_header_footer(slide6, "05 / 词汇复习计划与数据闭环流向", 6)
    
    # Left: Data Flow
    draw_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide6, "词汇羁绊闭环学习流 (Data Loop)", Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    flow_data = [
        {"text": "1. 学习触发 (Trigger)", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(12), "space_after": Pt(2)},
        {"text": "   用户在前端背词、答题或主动发起单词查询，激发交互数据流。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "2. 桥接同步 (Synchronization)", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   答题正确率、生词等数据瞬间通过 postMessage 传递至 C++ 后端进程。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "3. 物理持久化 (Persistence)", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   C++ 更新本地错题索引，计算单词的下一次复习时间戳并保存于本地 JSON。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(8)},
        
        {"text": "4. 桌宠陪伴反馈 (Feedback)", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(2)},
        {"text": "   根据成绩触发 C++ “伴读灵”动作更新与文字对话，增进词汇羁绊。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide6, Inches(1.1), Inches(2.4), Inches(5.0), Inches(3.8), flow_data)
    
    # Right: Priority Stack
    draw_card(slide6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide6, "智能复习优先级队列设计", Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    # We can design a visual list of 4 nested cards representing priorities!
    p_tops = [Inches(2.5), Inches(3.5), Inches(4.5), Inches(5.5)]
    p_colors = [RGBColor(210, 80, 80), COLOR_ACCENT, COLOR_PRIMARY, COLOR_TEXT_MUTED]
    p_titles = [
        "PRIORITY 01.  艾宾浩斯到期复习 (Overdue Review)",
        "PRIORITY 02.  用户主动查词计划 (Vocab Plans)",
        "PRIORITY 03.  高频错题专项巩固 (Incorrect Focus)",
        "PRIORITY 04.  日常设定今日新词 (Daily New Words)"
    ]
    p_descs = [
        "系统首先筛出记忆时间点已到期的词，强制用户复习，保证抗遗忘效果。",
        "处理用户在生词本或查词界面添加的待学词，保障即查即学流程。",
        "从历史答错词库中抽取进行高强度重新记忆与巩固测试。",
        "最后引入当天剩余的尚未接触过的新生词进行常规学习。"
    ]
    
    for idx in range(4):
        # Mini cards inside Priority list
        draw_card(slide6, Inches(7.1), p_tops[idx], Inches(5.2), Inches(0.85), COLOR_BG, p_colors[idx])
        # Colored bar
        draw_line(slide6, Inches(7.1), p_tops[idx], Inches(0.08), Inches(0.85), p_colors[idx])
        add_text(slide6, p_titles[idx], Inches(7.3), p_tops[idx] + Inches(0.08), Inches(4.9), Inches(0.3), Pt(11), p_colors[idx], bold=True)
        add_text(slide6, p_descs[idx], Inches(7.3), p_tops[idx] + Inches(0.38), Inches(4.9), Inches(0.4), Pt(9.5), COLOR_TEXT_DARK)

    # -------------------------------------------------------------
    # SLIDE 7: Build & Verification (系统运行与验证)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide7, COLOR_BG)
    add_header_footer(slide7, "06 / 系统的运行编译与集成测试验证", 7)
    
    # Left: Build workflow
    draw_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide7, "自动化构建与部署流", Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    build_steps = [
        {"text": "• CMake 跨平台管理", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(12), "space_after": Pt(3)},
        {"text": "   C++ 项目由 CMake 集中管理编译流程。配置中自动拉取并链接 Win32、Direct2D 和 Windows SAPI 底层库文件。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(12)},
        
        {"text": "• open_web.bat 一键批处理", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   编写了一键启动脚本，自动完成 CMake 编译检查、拉起本地网页服务，并初始化 C++ 悬浮桌宠进程，极大地简化了用户的运行配置难度。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(12)},
        
        {"text": "• 极低依赖，免安装发布", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   应用所有依赖项（Boost.JSON、WebView2 运行库等）均采用静态链接或本地打包，支持打包成绿色免安装文件夹直接双击运行。", "size": Pt(10.5), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide7, Inches(1.1), Inches(2.4), Inches(5.0), Inches(3.8), build_steps)
    
    # Right: Verification
    draw_card(slide7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide7, "三维立体集成测试验证", Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    test_steps = [
        {"text": "✓ 前端静态语法检查", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(12), "space_after": Pt(3)},
        {"text": "   通过 node --check web/app.js 保证前端主逻辑无 JavaScript 语法错误与未定义变量，维持前端高稳定性。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(12)},
        
        {"text": "✓ C++ 项目顺利编译", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   经由 MSVC 2022 及 GCC 环境完整测试，解决多线程死锁及 DLL 依赖，实现编译链接 0 警告、0 错误。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(12)},
        
        {"text": "✓ 功能级黑盒集成验证", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   系统运行后，拦截 /status 接口并校验心跳返回 \"ok\"；模拟连续答错，检测桌宠气泡能否在 4s 内弹出，保证人机伴学机制无时延延迟。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide7, Inches(7.2), Inches(2.4), Inches(5.0), Inches(3.8), test_steps)

    # -------------------------------------------------------------
    # SLIDE 8: Summary & Outlook (项目总结与展望)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide8, COLOR_BG)
    add_header_footer(slide8, "07 / 总结归纳与未来展望", 8)
    
    # Left: Project Summary
    draw_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide8, "《灵契·伴读灵》项目总结", Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    summary_data = [
        {"text": "• C++ 桌面宿主与 Web 前端深度融合", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(10), "space_after": Pt(3)},
        {"text": "   突破常规纯网页学习系统的局限性，将可交互的桌宠以 Win32 透明置顶窗口呈现，实现了真正的伴随式背词交互。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(10)},
        
        {"text": "• 系统底层多线程技术实践", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   深度整合 Win32 消息泵、Direct2D 硬件加速渲染、SAPI 多线程语音合成，让桌面程序交互零卡顿、发音极速响应。", "size": Pt(10.5), "color": COLOR_TEXT_DARK, "space_after": Pt(10)},
        
        {"text": "• 重视人机交互细节与健壮度", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   持续迭代布局细节，引入单词助记与历史纪录卡片，修复气泡截断、消除多余 HTML 桌宠，运行非常稳定。", "size": Pt(10.5), "color": COLOR_TEXT_DARK}
    ]
    add_rich_text(slide8, Inches(1.1), Inches(2.3), Inches(5.0), Inches(3.8), summary_data)
    
    # Right: Future Outlook
    draw_card(slide8, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.6), COLOR_WHITE, COLOR_BORDER)
    draw_card(slide8, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.5), COLOR_PRIMARY)
    add_text(slide8, "未来规划与优化方向", Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.3), Pt(12), COLOR_WHITE, bold=True)
    
    outlook_data = [
        {"text": "• 自适应间隔重复算法 (SM-2/Anki)", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_before": Pt(10), "space_after": Pt(3)},
        {"text": "   计划引入类似 SuperMemo-2 的主动间隔记忆算法，根据用户的掌握程度智能定制复习曲线，显著提升记忆成效。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(10)},
        
        {"text": "• 养成体系与社交互动扩展", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   计划丰富桌宠动作与皮肤库，引入“喂食”、“亲密度”、“羁绊等级”等玩法，激发用户的长期学习主动性。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED, "space_after": Pt(10)},
        
        {"text": "• 词库智能导入与分析报告", "size": Pt(13), "bold": True, "color": COLOR_PRIMARY, "space_after": Pt(3)},
        {"text": "   增加 PDF/Txt 智能提词功能，自动匹配用户专属词库，并提供阶段性词汇掌握度雷达图等数据图表分析。", "size": Pt(10.5), "color": COLOR_TEXT_MUTED}
    ]
    add_rich_text(slide8, Inches(7.2), Inches(2.3), Inches(5.0), Inches(3.8), outlook_data)

    # Save to file
    output_filename = "《灵契·伴读灵》--词汇羁绊养成系统答辩汇报.pptx"
    output_path = os.path.join("c:\\Users\\TYH\\Documents\\Codex\\2026-06-11\\c\\outputs\\word_trainer_web", output_filename)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

    # Also save a copy inside docs for clean backups
    docs_path = os.path.join("c:\\Users\\TYH\\Documents\\Codex\\2026-06-11\\c\\outputs\\word_trainer_web\\docs", output_filename)
    prs.save(docs_path)
    print(f"Backup presentation saved to: {docs_path}")

if __name__ == "__main__":
    create_presentation()
